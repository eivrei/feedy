from pptx import Presentation


# This method is run to extract information from .pptx files.
# It ignores all figures and pictures.
# Collects all topics and subpoints and returns a list on the form:
# [[title, point, point, ...], [title, point, point, ...], ...]
def extract(filepath):
    presentation = Presentation(filepath)
    all_data = []
    for slide in presentation.slides:
        slide_data = []
        number_of_sub_questions = 0
        for shape in slide.shapes:
            sub_points = []
            last_indent_level = 0  # Keeps track on what indent level the last point had.
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.level > 1:
                    continue
                paragraph_text = "".join(run.text for run in paragraph.runs)
                if paragraph_text != "":
                    # If point is subpoint -> make new question with subpoints as points
                    if paragraph.level > last_indent_level:
                        sub_points.extend([slide_data[-1], paragraph_text])
                        number_of_sub_questions += 1
                    elif paragraph.level == 1:
                        sub_points.append(paragraph_text)
                    elif paragraph.level < last_indent_level:
                        all_data.append(sub_points)
                        slide_data.append(paragraph_text)
                        sub_points = []
                    else:
                        slide_data.append(paragraph_text)
                    last_indent_level = paragraph.level

        # This line ensures that a subpoint with its points is behind its parent point in all_data.
        all_data.insert(len(all_data) - number_of_sub_questions, slide_data)
    return all_data

if __name__ == '__main__':
    print(extract("/Users/eivindreime/git/pugruppe100/server/temp/test_2.pptx"))